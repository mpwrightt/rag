"""
Utilities to convert various document types to Markdown for unified ingestion.
"""
from __future__ import annotations

import os
import io
import logging
import re
import unicodedata
from pathlib import Path
from typing import Tuple, Dict, Any

logger = logging.getLogger(__name__)

# Optional heavy deps imported lazily
try:
    from .docling_parser import get_docling_parser, is_docling_available
    DOCLING_AVAILABLE = is_docling_available()
except Exception:
    DOCLING_AVAILABLE = False
    logger.warning("Docling parser not available - falling back to traditional parsers")

def _looks_like_binary_text(text: str) -> bool:
    """Detect binary/gibberish text such as raw PDF bytes or ZIP/XML bodies.
    Returns True if content should be considered non-readable.
    """
    if not text:
        return False
    t = text.strip()
    # Common PDF/ZIP markers near the start
    if any(m in t[:300] for m in ("%PDF-", "startxref", "endobj", "PK\x03\x04")):
        return True
    # Heuristic: if fewer than 70% characters are letters/digits/basic punctuation
    import re as _re
    allowed = _re.compile(r"[\w\s\.,;:'\-\(\)\[\]/&%]", _re.UNICODE)
    total = len(t)
    if total >= 24:
        allowed_count = sum(1 for ch in t if allowed.match(ch))
        if allowed_count / max(1, total) < 0.7:
            return True
    # Extremely long unbroken tokens (compressed-looking)
    if any(len(tok) > 120 for tok in t.split()):
        return True
    return False

def _read_text(path: str, encodings=("utf-8", "latin-1")) -> str:
    for enc in encodings:
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except Exception:
            continue
    # Last resort: binary read then decode ignoring errors
    try:
        with open(path, "rb") as f:
            return f.read().decode("utf-8", errors="ignore")
    except Exception:
        return ""


def _normalize_text(text: str) -> str:
    """Normalize extracted text to improve downstream accuracy.
    - Unicode normalize (NFKC) to fix ligatures, fancy punctuation
    - Remove spurious control chars / replacement glyphs
    - Fix hyphenation across line breaks
    - Collapse excessive whitespace while preserving paragraphs
    """
    if not text:
        return ""
    t = unicodedata.normalize("NFKC", text)
    # Remove BOM and replacement characters
    t = t.replace("\ufeff", "").replace("�", "")
    # Normalize newlines
    t = t.replace("\r\n", "\n").replace("\r", "\n")
    # De-hyphenate line-break hyphenation: "com-\npany" -> "company"
    t = re.sub(r"-\s*\n\s*", "", t)
    # Merge lines that were wrapped mid-sentence: single newline between non-blank lines -> space
    t = re.sub(r"(?<=\S)\n(?=\S)", " ", t)
    # Collapse multiple blank lines to max two
    t = re.sub(r"\n{3,}", "\n\n", t)
    # Remove stray control chars except tab/newline
    t = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F]", "", t)
    # Collapse runs of spaces
    t = re.sub(r"[ \t]{2,}", " ", t)
    return t.strip()


def convert_to_markdown(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """
    Convert a file to Markdown text.

    Returns: (markdown_text, metadata)
    """
    p = Path(file_path)
    ext = p.suffix.lower().strip(".")
    meta: Dict[str, Any] = {
        "original_path": str(p),
        "original_ext": ext,
    }

    try:
        if ext in {"md", "markdown", "txt"}:
            text = _read_text(str(p))
            return text, meta

        if ext == "pdf":
            # First attempt: Docling (if available and enabled)
            if DOCLING_AVAILABLE and os.getenv("USE_DOCLING", "1").lower() in {"1", "true", "yes", "on"}:
                try:
                    docling_parser = get_docling_parser()
                    if docling_parser:
                        markdown_content, docling_meta = docling_parser.parse_document(str(p))
                        if markdown_content.strip():
                            logger.info(f"Successfully parsed {p.name} with Docling")
                            return _normalize_text(markdown_content), {**meta, **docling_meta}
                        else:
                            logger.warning("Docling parser returned empty content for %s", p.name)
                except Exception as e:
                    logger.warning("Docling parsing failed for %s: %s", p.name, e)
                    logger.info("Falling back to traditional PDF parsers")

            # Second attempt: pdfminer text extraction (fallback)
            try:
                from pdfminer.high_level import extract_text
                text = extract_text(str(p)) or ""
                if text.strip():
                    return _normalize_text(text), meta
                else:
                    logger.warning("PDF conversion via pdfminer returned empty text for %s", p.name)
            except Exception as e:
                logger.warning("PDF conversion failed via pdfminer: %s", e)

            # Third attempt: PyMuPDF (fitz) if available
            try:
                import fitz  # PyMuPDF
                try:
                    doc = fitz.open(str(p))
                    pages_text = []
                    for idx, page in enumerate(doc):
                        if idx >= 50:  # cap pages for very large PDFs
                            break
                        pages_text.append(page.get_text("text") or "")
                    doc.close()
                    text = "\n\n".join(pages_text)
                    if text.strip():
                        return _normalize_text(text), meta
                except Exception as e2:
                    logger.warning("PyMuPDF extraction failed: %s", e2)
            except Exception:
                # PyMuPDF not installed
                pass

            # Optional OCR fallback if enabled
            if os.getenv("OCR_PDF", "0").lower() in {"1", "true", "yes", "on"}:
                try:
                    from pdf2image import convert_from_path
                    import pytesseract
                    from PIL import Image
                except Exception as e:
                    logger.warning("OCR requested but pdf2image/pytesseract/Pillow not available: %s", e)
                else:
                    try:
                        poppler_path = os.getenv("POPPLER_PATH") or None
                        # Convert up to first 20 pages to limit processing time
                        images = convert_from_path(str(p), dpi=300, fmt="jpeg", poppler_path=poppler_path)
                        ocr_pages = []
                        for idx, img in enumerate(images):
                            if idx >= 20:
                                break
                            if not isinstance(img, Image.Image):
                                img = img.convert("RGB")
                            txt = pytesseract.image_to_string(img)
                            if txt and txt.strip():
                                ocr_pages.append(txt)
                        if ocr_pages:
                            return _normalize_text("\n\n".join(ocr_pages)), {**meta, "note": "ocr_fallback"}
                        else:
                            logger.warning("OCR produced no text for %s", p.name)
                    except Exception as e:
                        logger.warning("OCR fallback failed: %s", e)

            # Final fallback: try raw decode but drop if binary
            _raw = _read_text(str(p))
            if _looks_like_binary_text(_raw):
                logger.warning("PDF extraction failed and looked binary for %s; dropping text", p.name)
                return "", {**meta, "warning": "pdf extraction failed; no readable text"}
            return _normalize_text(_raw), {**meta, "warning": "pdf extraction failed; raw decode used"}

        if ext in {"docx"}:
            # Try Docling first (if available)
            if DOCLING_AVAILABLE and os.getenv("USE_DOCLING", "1").lower() in {"1", "true", "yes", "on"}:
                try:
                    docling_parser = get_docling_parser()
                    if docling_parser:
                        markdown_content, docling_meta = docling_parser.parse_document(str(p))
                        if markdown_content.strip():
                            logger.info(f"Successfully parsed {p.name} with Docling")
                            return _normalize_text(markdown_content), {**meta, **docling_meta}
                except Exception as e:
                    logger.warning("Docling parsing failed for %s: %s", p.name, e)

            # Fallback to python-docx
            try:
                from docx import Document
                doc = Document(str(p))
                parts = []
                for para in doc.paragraphs:
                    parts.append(para.text)
                text = "\n\n".join(filter(None, parts))
                return text, meta
            except Exception as e:
                logger.warning("DOCX conversion failed: %s", e)
                # Don't return raw ZIP bytes
                return "", {**meta, "warning": "docx read failed; no readable text"}

        if ext in {"pptx"}:
            # Try Docling first (if available)
            if DOCLING_AVAILABLE and os.getenv("USE_DOCLING", "1").lower() in {"1", "true", "yes", "on"}:
                try:
                    docling_parser = get_docling_parser()
                    if docling_parser:
                        markdown_content, docling_meta = docling_parser.parse_document(str(p))
                        if markdown_content.strip():
                            logger.info(f"Successfully parsed {p.name} with Docling")
                            return _normalize_text(markdown_content), {**meta, **docling_meta}
                except Exception as e:
                    logger.warning("Docling parsing failed for %s: %s", p.name, e)

            # Fallback to python-pptx
            try:
                from pptx import Presentation
                prs = Presentation(str(p))
                slides_text = []
                for s in prs.slides:
                    buf = []
                    for shape in s.shapes:
                        if hasattr(shape, "text"):
                            t = shape.text.strip()
                            if t:
                                buf.append(t)
                    if buf:
                        slides_text.append("\n".join(buf))
                return "\n\n---\n\n".join(slides_text), meta
            except Exception as e:
                logger.warning("PPTX conversion failed: %s", e)
                # Don't return raw ZIP bytes
                return "", {**meta, "warning": "pptx read failed; no readable text"}

        if ext in {"html", "htm"}:
            try:
                from bs4 import BeautifulSoup
                from markdownify import markdownify as md
                html = _read_text(str(p))
                soup = BeautifulSoup(html, "html.parser")
                # Remove script/style
                for tag in soup(["script", "style"]):
                    tag.decompose()
                markdown = md(str(soup), heading_style="ATX")
                return markdown, meta
            except Exception as e:
                logger.warning("HTML conversion failed: %s", e)
                return _read_text(str(p)), {**meta, "warning": "html to md failed"}

        if ext in {"csv", "tsv"}:
            try:
                import csv
                delimiter = "," if ext == "csv" else "\t"
                out = io.StringIO()
                writer = out.write
                with open(str(p), newline="", encoding="utf-8") as f:
                    reader = csv.reader(f, delimiter=delimiter)
                    rows = list(reader)
                if rows:
                    # Simple Markdown table
                    header = rows[0]
                    writer("| " + " | ".join(header) + " |\n")
                    writer("|" + "|".join([" --- "] * len(header)) + "|\n")
                    for r in rows[1:]:
                        writer("| " + " | ".join(r) + " |\n")
                return out.getvalue(), meta
            except Exception as e:
                logger.warning("CSV/TSV conversion failed: %s", e)
                return _read_text(str(p)), {**meta, "warning": "csv/tsv to md failed"}

        if ext in {"xlsx", "xls"}:
            try:
                import pandas as pd
                xls = pd.ExcelFile(str(p))
                md_parts = []
                for sheet in xls.sheet_names:
                    df = xls.parse(sheet)
                    md_parts.append(f"\n\n## {sheet}\n\n")
                    md_parts.append(df.to_markdown(index=False))
                return "".join(md_parts), meta
            except Exception as e:
                logger.warning("Excel conversion via pandas failed: %s", e)
                # Fallback: openpyxl direct parsing (no pandas)
                try:
                    from openpyxl import load_workbook
                    wb = load_workbook(filename=str(p), data_only=True)
                    md_parts = []
                    for ws in wb.worksheets:
                        md_parts.append(f"\n\n## {ws.title}\n\n")
                        # Collect rows
                        rows = []
                        for row in ws.iter_rows(values_only=True):
                            rows.append(["" if v is None else str(v) for v in row])
                        if rows:
                            # Build simple Markdown table
                            headers = rows[0]
                            md_parts.append("| " + " | ".join(headers) + " |\n")
                            md_parts.append("|" + "|".join([" --- "] * len(headers)) + "|\n")
                            for r in rows[1:]:
                                md_parts.append("| " + " | ".join(r) + " |\n")
                    return "".join(md_parts), {**meta, "note": "excel parsed via openpyxl fallback"}
                except Exception as e2:
                    logger.warning("Excel conversion via openpyxl failed: %s", e2)
                    # Don't return raw binary
                    return "", {**meta, "warning": "excel to md failed; no readable text"}

        # Images and others: best-effort (no OCR by default)
        if ext in {"png", "jpg", "jpeg", "gif", "webp", "svg"}:
            # Placeholder markdown with link
            return f"![{p.name}]({p.name})", {**meta, "note": "image placeholder; OCR not enabled"}

        # Default fallback: attempt text read
        return _read_text(str(p)), {**meta, "warning": "unknown extension; raw text read"}

    except Exception as e:
        logger.error("Failed to convert %s: %s", file_path, e)
        return "", {**meta, "error": str(e)}
